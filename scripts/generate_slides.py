"""Generate demo video slides as PPTX for DocuAlign AI."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
import os
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Constants ────────────────────────────────────────────────────
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)       # Deep navy
ACCENT_GREEN = RGBColor(0x34, 0xD3, 0x99)  # Emerald green
ACCENT_BLUE = RGBColor(0x38, 0xBD, 0xF8)   # Sky blue
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xA0, 0xA0, 0xB0)
CARD_BG = RGBColor(0x1E, 0x29, 0x3B)       # Dark card
CRITICAL_RED = RGBColor(0xEF, 0x44, 0x44)
WARNING_YELLOW = RGBColor(0xFB, 0xBF, 0x24)

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_DIR = os.path.dirname(SCRIPT_DIR)
LOGO_PATH = os.path.join(PROJECT_DIR, "assets", "logo.png")
ARCH_RUNTIME_PATH = os.path.join(PROJECT_DIR, "docs", "architecture_runtime.png")
ARCH_DEV_PATH = os.path.join(PROJECT_DIR, "docs", "architecture_development.png")
INFO_BLUE = RGBColor(0x60, 0xA5, 0xFA)

SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color,
                     text="", font_size=14, text_color=WHITE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
    return shape


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # Blank

    # =====================================================================
    # SLIDE 1: Title Card
    # =====================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1, DARK_BG)

    # Accent bar at top
    bar = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_GREEN
    bar.line.fill.background()

    # Logo image
    if os.path.exists(LOGO_PATH):
        slide1.shapes.add_picture(
            LOGO_PATH, Inches(9.0), Inches(1.2), Inches(3.5), Inches(3.5)
        )
    else:
        circle = slide1.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(9.0), Inches(1.2), Inches(3.5), Inches(3.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT_GREEN
        circle.line.fill.background()


    # Main title
    add_text_box(slide1, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
                 "DocuAlign AI", font_size=60, color=WHITE, bold=True)

    # Subtitle
    add_text_box(slide1, Inches(1), Inches(3.3), Inches(8), Inches(0.8),
                 "AI-Powered Document Integrity Agent", font_size=28,
                 color=ACCENT_GREEN)

    # Description
    add_text_box(slide1, Inches(1), Inches(4.5), Inches(10), Inches(1.0),
                 "Gemini 2.0 Flash × LangGraph × Google Cloud",
                 font_size=20, color=LIGHT_GRAY)

    # Bottom tagline
    add_text_box(slide1, Inches(1), Inches(6.0), Inches(11), Inches(0.6),
                 "ドキュメントの「サイレント劣化」をAIが自動検知", font_size=18,
                 color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

    # =====================================================================
    # SLIDE 2: Problem Statement
    # =====================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide2, DARK_BG)

    add_text_box(slide2, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                 "大企業が抱える「ドキュメント劣化」問題", font_size=36,
                 color=WHITE, bold=True)

    add_text_box(slide2, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
                 "既存ツールでは解決できない — 差分ではなく「意味的な矛盾」を検出する技術が存在しなかった",
                 font_size=16, color=LIGHT_GRAY)

    # 3 stat cards (industry research-backed)
    stats = [
        ("21.3%",
         "の生産性が\nドキュメント管理の\n非効率で失われている",
         "📉", CRITICAL_RED,
         "Iron Mountain 調査"),
        ("$19,732",
         "/人・年のコストが\n情報検索・文書管理に\n費やされている",
         "💰", WARNING_YELLOW,
         "IDC / Ripcord 調査"),
        ("2.5h/日",
         "を従業員が\n必要な情報の検索に\n費やしている",
         "⏱️", INFO_BLUE,
         "Forbes / McKinsey"),
    ]

    for i, (number, desc, icon, accent, source) in enumerate(stats):
        x = Inches(0.8 + i * 4.0)
        y = Inches(2.1)

        # Card
        card = add_rounded_rect(slide2, x, y, Inches(3.5), Inches(4.5), CARD_BG)

        # Icon
        add_text_box(slide2, x + Inches(0.3), y + Inches(0.3),
                     Inches(1), Inches(0.8), icon, font_size=40)

        # Big number
        add_text_box(slide2, x + Inches(0.3), y + Inches(1.2),
                     Inches(3), Inches(1.0), number, font_size=44,
                     color=accent, bold=True)

        # Description
        add_text_box(slide2, x + Inches(0.3), y + Inches(2.5),
                     Inches(3), Inches(1.2), desc, font_size=16,
                     color=LIGHT_GRAY)

        # Source
        add_text_box(slide2, x + Inches(0.3), y + Inches(3.8),
                     Inches(3), Inches(0.5), source, font_size=11,
                     color=RGBColor(0x70, 0x70, 0x80))

    # =====================================================================
    # SLIDE 3: Runtime Architecture (editable shapes)
    # =====================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide3, DARK_BG)

    add_text_box(slide3, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "Runtime Architecture", font_size=34,
                 color=WHITE, bold=True)
    add_text_box(slide3, Inches(0.8), Inches(0.9), Inches(11), Inches(0.35),
                 "100% Google Cloud — Serverless & Fully Managed",
                 font_size=15, color=ACCENT_GREEN)

    # --- Row 1: User → Cloud Run → LangGraph Agent ---
    row1_boxes = [
        ("👤 ユーザー", "ブラウザ", ACCENT_BLUE, Inches(0.3)),
        ("Cloud Run", "Streamlit App", ACCENT_GREEN, Inches(3.0)),
        ("LangGraph", "自律型Agent", ACCENT_GREEN, Inches(5.7)),
        ("Vertex AI", "Gemini 2.0 Flash", ACCENT_GREEN, Inches(8.4)),
        ("Firestore", "結果保存", ACCENT_BLUE, Inches(11.1)),
    ]
    for name, desc, accent, x in row1_boxes:
        add_rounded_rect(slide3, x, Inches(1.5), Inches(2.3), Inches(1.4), CARD_BG)
        add_text_box(slide3, x + Inches(0.1), Inches(1.55),
                     Inches(2.1), Inches(0.55), name, font_size=14,
                     color=accent, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide3, x + Inches(0.1), Inches(2.15),
                     Inches(2.1), Inches(0.55), desc, font_size=11,
                     color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Arrows row 1
    for ax in [Inches(2.6), Inches(5.3), Inches(8.0), Inches(10.7)]:
        add_text_box(slide3, ax, Inches(1.85), Inches(0.4), Inches(0.4),
                     "→", font_size=22, color=ACCENT_GREEN,
                     alignment=PP_ALIGN.CENTER)

    # --- Row 2: Data Sources (left) ---
    add_text_box(slide3, Inches(0.3), Inches(3.15), Inches(3), Inches(0.4),
                 "📄 データソース", font_size=14, color=ACCENT_BLUE, bold=True)
    data_sources = [
        ("Google Drive", "ドキュメント取得"),
        ("Cloud Storage", "ファイル保存"),
        ("Google Docs", "リアルタイム同期"),
    ]
    for i, (name, desc) in enumerate(data_sources):
        y = Inches(3.6 + i * 0.85)
        add_rounded_rect(slide3, Inches(0.3), y, Inches(2.8), Inches(0.75), CARD_BG)
        add_text_box(slide3, Inches(0.4), y + Inches(0.05),
                     Inches(2.6), Inches(0.35), name, font_size=12,
                     color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide3, Inches(0.4), y + Inches(0.38),
                     Inches(2.6), Inches(0.3), desc, font_size=10,
                     color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # --- Row 2: Infrastructure (center) ---
    add_text_box(slide3, Inches(3.5), Inches(3.15), Inches(5), Inches(0.4),
                 "⚙️ インフラストラクチャ", font_size=14,
                 color=ACCENT_BLUE, bold=True)
    infra = [
        ("Cloud Tasks", "非同期処理"),
        ("Pub/Sub", "リアルタイム通知"),
        ("Eventarc", "イベントトリガー"),
    ]
    for i, (name, desc) in enumerate(infra):
        x = Inches(3.5 + i * 2.1)
        add_rounded_rect(slide3, x, Inches(3.6), Inches(1.9), Inches(0.75), CARD_BG)
        add_text_box(slide3, x + Inches(0.1), Inches(3.65),
                     Inches(1.7), Inches(0.35), name, font_size=12,
                     color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide3, x + Inches(0.1), Inches(3.98),
                     Inches(1.7), Inches(0.3), desc, font_size=10,
                     color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # --- Row 2: Security & Ops (right) ---
    add_text_box(slide3, Inches(10.0), Inches(3.15), Inches(3), Inches(0.4),
                 "🔒 セキュリティ & 運用", font_size=14,
                 color=ACCENT_BLUE, bold=True)
    sec_ops = [
        ("Secret Manager", "機密情報管理"),
        ("Cloud Logging", "監査ログ"),
        ("IAM", "アクセス制御"),
    ]
    for i, (name, desc) in enumerate(sec_ops):
        y = Inches(3.6 + i * 0.85)
        add_rounded_rect(slide3, Inches(10.0), y, Inches(2.8), Inches(0.75), CARD_BG)
        add_text_box(slide3, Inches(10.1), y + Inches(0.05),
                     Inches(2.6), Inches(0.35), name, font_size=12,
                     color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide3, Inches(10.1), y + Inches(0.38),
                     Inches(2.6), Inches(0.3), desc, font_size=10,
                     color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Differentiator bar
    add_rounded_rect(slide3, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.65),
                     ACCENT_GREEN)
    add_text_box(slide3, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.55),
                 "💡 差別化: テキスト意味矛盾 + 画像劣化のマルチモーダル検出 — "
                 "LangGraph エージェントが自律実行",
                 font_size=15, color=DARK_BG, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # =====================================================================
    # SLIDE 4: Development Architecture / Google AntiGravity (editable)
    # =====================================================================
    slide4_dev = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide4_dev, DARK_BG)

    add_text_box(slide4_dev, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "Development with Google AntiGravity",
                 font_size=34, color=WHITE, bold=True)
    add_text_box(slide4_dev, Inches(0.8), Inches(0.9), Inches(11), Inches(0.35),
                 "AI-Assisted Coding — 設計から本番デプロイまで一気通貫",
                 font_size=15, color=ACCENT_GREEN)

    # --- Dev Pipeline Flow (top row) ---
    dev_flow = [
        ("🧑‍💻 Developer", "設計 & コーディング", ACCENT_BLUE),
        ("AntiGravity", "AI アシスト\nコード生成", ACCENT_GREEN),
        ("GitHub", "バージョン管理\nPR & Code Review", ACCENT_BLUE),
        ("Cloud Build", "CI/CD\n自動ビルド & テスト", ACCENT_BLUE),
        ("Cloud Run", "本番デプロイ\n自動スケーリング", ACCENT_GREEN),
    ]
    for i, (name, desc, accent) in enumerate(dev_flow):
        x = Inches(0.2 + i * 2.6)
        add_rounded_rect(slide4_dev, x, Inches(1.5), Inches(2.3), Inches(1.6),
                         CARD_BG)
        add_text_box(slide4_dev, x + Inches(0.1), Inches(1.55),
                     Inches(2.1), Inches(0.5), name, font_size=14,
                     color=accent, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide4_dev, x + Inches(0.1), Inches(2.15),
                     Inches(2.1), Inches(0.75), desc, font_size=11,
                     color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
        if i < len(dev_flow) - 1:
            add_text_box(slide4_dev, x + Inches(2.3), Inches(2.0),
                         Inches(0.3), Inches(0.4), "→", font_size=22,
                         color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

    # --- AntiGravity contribution areas (bottom) ---
    add_text_box(slide4_dev, Inches(0.5), Inches(3.4), Inches(12), Inches(0.4),
                 "🛠️ AntiGravity が支援した開発領域",
                 font_size=16, color=ACCENT_GREEN, bold=True)

    contributions = [
        ("Architecture\nDesign", "GCPサービス構成\n最適な設計を提案",
         "設計時間 60%削減", ACCENT_BLUE),
        ("LangGraph\nPipeline", "エージェント\nパイプライン実装",
         "実装速度 3倍", ACCENT_GREEN),
        ("Streamlit\nDashboard", "UI/UXデザイン\nコンポーネント生成",
         "UI開発 70%削減", ACCENT_BLUE),
        ("Unit & E2E\nTests", "テストケース\n自動生成・実行",
         "テスト工数 80%削減", ACCENT_GREEN),
        ("Documentation\n& Slides", "コード・プレゼン\n資料の自動生成",
         "ドキュメント 90%自動", ACCENT_BLUE),
    ]

    for i, (title, desc, metric, accent) in enumerate(contributions):
        x = Inches(0.2 + i * 2.6)
        y = Inches(3.9)
        add_rounded_rect(slide4_dev, x, y, Inches(2.3), Inches(2.6), CARD_BG)
        add_text_box(slide4_dev, x + Inches(0.1), y + Inches(0.1),
                     Inches(2.1), Inches(0.7), title, font_size=13,
                     color=accent, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide4_dev, x + Inches(0.1), y + Inches(0.85),
                     Inches(2.1), Inches(0.8), desc, font_size=10,
                     color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
        # Metric badge
        add_rounded_rect(slide4_dev, x + Inches(0.15), y + Inches(1.9),
                         Inches(2.0), Inches(0.5), accent)
        add_text_box(slide4_dev, x + Inches(0.15), y + Inches(1.93),
                     Inches(2.0), Inches(0.45), metric, font_size=11,
                     color=DARK_BG, bold=True, alignment=PP_ALIGN.CENTER)

    # Bottom bar
    add_rounded_rect(slide4_dev, Inches(0.3), Inches(7.0), Inches(12.7),
                     Inches(0.25), ACCENT_GREEN)

    # =====================================================================
    # SLIDE 5: Architecture & Core Tech (pipeline detail)
    # =====================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide3, DARK_BG)

    add_text_box(slide3, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                 "🤖 アーキテクチャ & コア技術", font_size=36,
                 color=WHITE, bold=True)

    # Pipeline flow
    steps = [
        ("1. 取得", "Google Drive\nから文書取得", ACCENT_BLUE),
        ("2. 検索", "Agent Builder\nで関連文書検索", ACCENT_BLUE),
        ("3. 分析", "Gemini 2.0\nテキスト矛盾検出", ACCENT_GREEN),
        ("4. 画像", "Gemini 2.0\nスクショ劣化検出", ACCENT_GREEN),
        ("5. 提案", "修正提案\n自動生成", ACCENT_BLUE),
    ]

    for i, (title, desc, color) in enumerate(steps):
        x = Inches(0.5 + i * 2.5)
        y = Inches(1.8)

        box = add_rounded_rect(slide3, x, y, Inches(2.2), Inches(2.0), CARD_BG)
        add_text_box(slide3, x + Inches(0.15), y + Inches(0.2),
                     Inches(1.9), Inches(0.5), title, font_size=16,
                     color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide3, x + Inches(0.15), y + Inches(0.8),
                     Inches(1.9), Inches(1.0), desc, font_size=13,
                     color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

        # Arrow between steps
        if i < len(steps) - 1:
            add_text_box(slide3, x + Inches(2.2), y + Inches(0.7),
                         Inches(0.3), Inches(0.5), "→", font_size=24,
                         color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

    # GCP Services row
    services = [
        ("Cloud Run", "サーバーレス\nホスティング"),
        ("Firestore", "スキャン結果\n保存"),
        ("GCS", "ドキュメント\nストレージ"),
        ("Eventarc", "イベント\nトリガー"),
        ("Pub/Sub", "リアルタイム\n通知"),
    ]

    add_text_box(slide3, Inches(0.8), Inches(4.2), Inches(5), Inches(0.5),
                 "☁️ Google Cloud サービス", font_size=18,
                 color=ACCENT_BLUE, bold=True)

    for i, (name, desc) in enumerate(services):
        x = Inches(0.5 + i * 2.5)
        y = Inches(4.8)
        box = add_rounded_rect(slide3, x, y, Inches(2.2), Inches(1.5), CARD_BG)
        add_text_box(slide3, x + Inches(0.15), y + Inches(0.15),
                     Inches(1.9), Inches(0.4), name, font_size=15,
                     color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide3, x + Inches(0.15), y + Inches(0.6),
                     Inches(1.9), Inches(0.8), desc, font_size=12,
                     color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Key differentiator
    diff_box = add_rounded_rect(
        slide3, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7),
        ACCENT_GREEN
    )
    add_text_box(slide3, Inches(0.8), Inches(6.55), Inches(11.8), Inches(0.6),
                 "💡 差別化: テキスト意味矛盾 + 画像劣化のマルチモーダル検出 — "
                 "LangGraphエージェントが自律実行",
                 font_size=16, color=DARK_BG, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # =====================================================================
    # SLIDE 6: Business Value / ROI (refined)
    # =====================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide4, DARK_BG)

    add_text_box(slide4, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "📈 導入効果 — Before / After", font_size=34,
                 color=WHITE, bold=True)
    add_text_box(slide4, Inches(0.8), Inches(0.9), Inches(11), Inches(0.35),
                 "50人規模の技術組織における年間試算（ドキュメント200件想定）",
                 font_size=14, color=LIGHT_GRAY)

    # Before/After comparison rows
    comparisons = [
        ("ドキュメント整合性チェック",
         "手動レビュー\n40h/月 × 12 = 480h/年",
         "AI自動スキャン\n2h/月 × 12 = 24h/年",
         "96%", "時間削減"),
        ("古い手順書による障害対応",
         "月3件 × 復旧4h\n= 144h/年",
         "月0.5件 × 復旧4h\n= 24h/年",
         "83%", "インシデント減"),
        ("新人オンボーディング遅延",
         "古い手順で平均\n+2日/人 × 20人 = 40日/年",
         "最新ドキュメント維持\n+0.5日/人 = 10日/年",
         "75%", "遅延削減"),
        ("コンプライアンス違反リスク",
         "手動監査\n年2回 × 80h = 160h",
         "常時監視 + アラート\n年2回 × 20h = 40h",
         "75%", "監査工数減"),
    ]

    # Column headers
    header_y = Inches(1.4)
    add_text_box(slide4, Inches(0.4), header_y, Inches(2.8), Inches(0.4),
                 "項目", font_size=13, color=ACCENT_BLUE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_rounded_rect(slide4, Inches(3.3), header_y, Inches(3.2), Inches(0.4),
                     CRITICAL_RED)
    add_text_box(slide4, Inches(3.3), header_y, Inches(3.2), Inches(0.4),
                 "Before（従来）", font_size=13, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_rounded_rect(slide4, Inches(6.6), header_y, Inches(3.2), Inches(0.4),
                     ACCENT_GREEN)
    add_text_box(slide4, Inches(6.6), header_y, Inches(3.2), Inches(0.4),
                 "After（DocuAlign AI）", font_size=13,
                 color=DARK_BG, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide4, Inches(10.0), header_y, Inches(3.0), Inches(0.4),
                 "削減効果", font_size=13, color=ACCENT_GREEN, bold=True,
                 alignment=PP_ALIGN.CENTER)

    for i, (item, before, after, pct, label) in enumerate(comparisons):
        y = Inches(1.95 + i * 1.2)

        # Row background (alternating)
        if i % 2 == 0:
            add_rounded_rect(slide4, Inches(0.3), y - Inches(0.05),
                             Inches(12.7), Inches(1.1),
                             RGBColor(0x15, 0x1F, 0x32))

        # Item name
        add_text_box(slide4, Inches(0.4), y, Inches(2.8), Inches(1.0),
                     item, font_size=13, color=WHITE, bold=True)

        # Before
        add_text_box(slide4, Inches(3.3), y, Inches(3.2), Inches(1.0),
                     before, font_size=11, color=CRITICAL_RED)

        # After
        add_text_box(slide4, Inches(6.6), y, Inches(3.2), Inches(1.0),
                     after, font_size=11, color=ACCENT_GREEN)

        # Percentage badge
        add_rounded_rect(slide4, Inches(10.3), y + Inches(0.05),
                         Inches(1.3), Inches(0.5), ACCENT_GREEN)
        add_text_box(slide4, Inches(10.3), y + Inches(0.05),
                     Inches(1.3), Inches(0.5), pct, font_size=20,
                     color=DARK_BG, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide4, Inches(10.3), y + Inches(0.55),
                     Inches(2.2), Inches(0.35), label, font_size=11,
                     color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

    # Bottom summary bar
    summary_y = Inches(6.8)
    add_rounded_rect(slide4, Inches(0.3), Inches(6.5), Inches(12.7),
                     Inches(0.8), CARD_BG)
    add_text_box(slide4, Inches(0.5), Inches(6.55), Inches(5), Inches(0.7),
                 "💰 年間コスト削減効果（人件費単価 ¥6,000/h 試算）",
                 font_size=14, color=WHITE, bold=True)
    add_text_box(slide4, Inches(7.0), Inches(6.5), Inches(6.0), Inches(0.8),
                 "約480万円 / 年　（800h × ¥6,000）",
                 font_size=28, color=ACCENT_GREEN, bold=True,
                 alignment=PP_ALIGN.RIGHT)

    # =====================================================================
    # SLIDE 5: Closing
    # =====================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide5, DARK_BG)

    # Accent bar
    bar = slide5.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_GREEN
    bar.line.fill.background()

    # Logo image on closing slide
    if os.path.exists(LOGO_PATH):
        slide5.shapes.add_picture(
            LOGO_PATH, Inches(5.4), Inches(0.5), Inches(2.5), Inches(2.5)
        )

    add_text_box(slide5, Inches(1), Inches(3.2), Inches(11), Inches(1.2),
                 "DocuAlign AI", font_size=60, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide5, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
                 "ドキュメントの「サイレント劣化」を\nAIが自動で検知・修正提案する世界へ",
                 font_size=24, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # GitHub card
    gh_box = add_rounded_rect(
        slide5, Inches(3), Inches(4.0), Inches(7.3), Inches(1.6), CARD_BG
    )
    add_text_box(slide5, Inches(3.5), Inches(4.15), Inches(6.3), Inches(0.5),
                 "🔗 GitHub Repository", font_size=18, color=ACCENT_BLUE,
                 bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide5, Inches(3.5), Inches(4.7), Inches(6.3), Inches(0.7),
                 "github.com/Koki0812/docugardener-agent",
                 font_size=22, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

    # Thank you
    add_text_box(slide5, Inches(1), Inches(6.2), Inches(11), Inches(0.6),
                 "ありがとうございました", font_size=28,
                 color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    return prs


if __name__ == "__main__":
    prs = build_presentation()
    output_path = "docs/demo_slides.pptx"
    prs.save(output_path)
    print(f"Saved: {output_path}")
